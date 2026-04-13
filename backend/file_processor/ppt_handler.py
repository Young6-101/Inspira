"""
PPT handler — extracts text AND embedded images from .pptx files.

Returns a structured result so the caller can:
  - text chunks → CLIP text encode → text_collection
  - image bytes → CLIP image encode → image_collection
"""
from dataclasses import dataclass, field
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class PPTExtractResult:
	text: str = ""
	images: list[dict] = field(default_factory=list)
	# images: [{"bytes": b"...", "name": "slide3_img1.png"}, ...]


def extract_from_pptx(file_path: str) -> PPTExtractResult:
	"""Extract text + embedded images from a .pptx file.

	Returns:
		PPTExtractResult with combined text and list of image blobs.
	"""
	path = Path(file_path)
	if not path.exists():
		print(f"--- [ERROR] PPTX file not found: {file_path} ---")
		return PPTExtractResult()

	try:
		prs = Presentation(file_path)
	except Exception as e:
		print(f"--- [ERROR] Failed to open PPTX {file_path}: {e} ---")
		return PPTExtractResult()

	all_parts: list[str] = []
	images: list[dict] = []
	img_counter = 0

	for slide_idx, slide in enumerate(prs.slides, start=1):
		slide_parts: list[str] = [f"[Slide {slide_idx}]"]

		for shape in slide.shapes:
			# --- Text extraction ---
			if shape.has_text_frame:
				for paragraph in shape.text_frame.paragraphs:
					text = paragraph.text.strip()
					if text:
						slide_parts.append(text)

			# --- Table extraction ---
			if shape.has_table:
				table = shape.table
				for row in table.rows:
					row_text = " | ".join(
						cell.text.strip() for cell in row.cells
					)
					if row_text.strip("| "):
						slide_parts.append(row_text)

			# --- Image extraction ---
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				try:
					image = shape.image
					image_bytes = image.blob
					ext = image.content_type.split("/")[-1]
					img_counter += 1
					img_name = f"slide{slide_idx}_img{img_counter}.{ext}"
					images.append({"bytes": image_bytes, "name": img_name})
					slide_parts.append(f"[Embedded image: {img_name}]")
				except Exception as e:
					print(f"--- [WARN] Could not extract image on slide {slide_idx}: {e} ---")

		slide_text = "\n".join(slide_parts)
		if slide_text.strip():
			all_parts.append(slide_text)

	result_text = "\n\n".join(all_parts)
	clean_text = " ".join(result_text.split())
	print(f"--- [LOG] Extracted {len(prs.slides)} slides, {len(images)} images from {path.name} ---")

	return PPTExtractResult(text=clean_text, images=images)


# Backward compat wrapper
def extract_text_from_pptx(file_path: str, describe_images: bool = False) -> str:
	"""Legacy API — returns text only."""
	return extract_from_pptx(file_path).text
